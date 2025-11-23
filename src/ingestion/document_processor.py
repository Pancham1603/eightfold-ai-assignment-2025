"""
Document Ingestion Pipeline for Eightfold AI Reference Materials
Processes various document formats and adds them to Pinecone vector store
"""

import os
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime
import hashlib

# Document processing libraries
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

from langchain_core.documents import Document
try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    from langchain.text_splitter import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Process and ingest Eightfold AI reference documents into vector store"""
    
    SUPPORTED_EXTENSIONS = {
        '.pdf': 'PDF Document',
        '.docx': 'Word Document',
        '.doc': 'Word Document',
        '.pptx': 'PowerPoint Presentation',
        '.ppt': 'PowerPoint Presentation',
        '.txt': 'Text File',
        '.md': 'Markdown File',
        '.xlsx': 'Excel Spreadsheet',
        '.xls': 'Excel Spreadsheet',
    }
    
    def __init__(self, vector_store, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize document processor
        
        Args:
            vector_store: Pinecone vector store instance
            chunk_size: Size of text chunks for embedding
            chunk_overlap: Overlap between chunks
        """
        self.vector_store = vector_store
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        
        # Check available libraries
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Check which document processing libraries are available"""
        missing = []
        
        if PdfReader is None:
            missing.append("pypdf (for PDF processing)")
        if DocxDocument is None:
            missing.append("python-docx (for Word documents)")
        if Presentation is None:
            missing.append("python-pptx (for PowerPoint)")
        if openpyxl is None:
            missing.append("openpyxl (for Excel files)")
        
        if missing:
            logger.warning(f"Missing optional dependencies: {', '.join(missing)}")
            logger.warning("Install with: pip install pypdf python-docx python-pptx openpyxl")
    
    def process_folder(
        self, 
        folder_path: str, 
        document_type: str = "eightfold_reference",
        metadata: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        Process all supported documents in a folder and add to vector store
        
        Args:
            folder_path: Path to folder containing Eightfold AI documents
            document_type: Type of documents (e.g., 'eightfold_reference', 'product_docs')
            metadata: Additional metadata to attach to all documents
        
        Returns:
            Dictionary with processing statistics
        """
        folder_path = Path(folder_path)
        
        if not folder_path.exists():
            raise ValueError(f"Folder does not exist: {folder_path}")
        
        if not folder_path.is_dir():
            raise ValueError(f"Path is not a directory: {folder_path}")
        
        stats = {
            'total_files': 0,
            'processed': 0,
            'failed': 0,
            'skipped': 0,
            'total_chunks': 0,
            'files_by_type': {},
            'errors': []
        }
        
        # Find all supported files
        supported_files = []
        for ext in self.SUPPORTED_EXTENSIONS.keys():
            supported_files.extend(folder_path.rglob(f"*{ext}"))
        
        stats['total_files'] = len(supported_files)
        logger.info(f"Found {len(supported_files)} supported documents in {folder_path}")
        
        # Process each file
        for file_path in supported_files:
            try:
                result = self.process_document(
                    str(file_path),
                    document_type=document_type,
                    metadata=metadata
                )
                
                if result['success']:
                    stats['processed'] += 1
                    stats['total_chunks'] += result['chunks_created']
                    
                    # Track by file type
                    ext = file_path.suffix.lower()
                    file_type = self.SUPPORTED_EXTENSIONS.get(ext, 'Unknown')
                    stats['files_by_type'][file_type] = stats['files_by_type'].get(file_type, 0) + 1
                    
                    logger.info(f"✓ Processed {file_path.name}: {result['chunks_created']} chunks")
                else:
                    stats['failed'] += 1
                    stats['errors'].append({
                        'file': str(file_path),
                        'error': result.get('error', 'Unknown error')
                    })
                    logger.error(f"✗ Failed to process {file_path.name}: {result.get('error')}")
                    
            except Exception as e:
                stats['failed'] += 1
                stats['errors'].append({
                    'file': str(file_path),
                    'error': str(e)
                })
                logger.error(f"✗ Error processing {file_path.name}: {e}")
        
        return stats
    
    def process_document(
        self,
        file_path: str,
        document_type: str = "eightfold_reference",
        metadata: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        Process a single document and add to vector store
        
        Args:
            file_path: Path to the document
            document_type: Type of document (for metadata tagging)
            metadata: Additional metadata
        
        Returns:
            Dictionary with processing results
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            return {'success': False, 'error': 'File not found'}
        
        # Extract text based on file type
        extension = file_path.suffix.lower()
        
        try:
            if extension == '.pdf':
                text = self._extract_pdf(file_path)
            elif extension in ['.docx', '.doc']:
                text = self._extract_docx(file_path)
            elif extension in ['.pptx', '.ppt']:
                text = self._extract_pptx(file_path)
            elif extension in ['.txt', '.md']:
                text = self._extract_text(file_path)
            elif extension in ['.xlsx', '.xls']:
                text = self._extract_excel(file_path)
            else:
                return {'success': False, 'error': f'Unsupported file type: {extension}'}
            
            if not text or len(text.strip()) < 10:
                return {'success': False, 'error': 'No text content extracted'}
            
            # Create document chunks
            chunks = self.text_splitter.split_text(text)
            
            # Prepare metadata
            doc_metadata = {
                'source': str(file_path),
                'filename': file_path.name,
                'file_type': self.SUPPORTED_EXTENSIONS.get(extension, 'Unknown'),
                'document_type': document_type,
                'ingestion_date': datetime.now().isoformat(),
                'chunk_count': len(chunks),
                'file_hash': self._compute_file_hash(file_path),
                'is_eightfold_reference': document_type == 'eightfold_reference'
            }
            
            # Add custom metadata
            if metadata:
                doc_metadata.update(metadata)
            
            # Create LangChain documents
            documents = []
            for i, chunk in enumerate(chunks):
                chunk_metadata = doc_metadata.copy()
                chunk_metadata['chunk_index'] = i
                
                documents.append(Document(
                    page_content=chunk,
                    metadata=chunk_metadata
                ))
            
            # Add to vector store with Eightfold AI context
            self.vector_store.add_eightfold_documents(documents)
            
            return {
                'success': True,
                'chunks_created': len(chunks),
                'file': str(file_path),
                'metadata': doc_metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            return {'success': False, 'error': str(e)}
    
    def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF"""
        if PdfReader is None:
            raise ImportError("pypdf not installed. Install with: pip install pypdf")
        
        reader = PdfReader(str(file_path))
        text = []
        
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
        
        return "\n\n".join(text)
    
    def _extract_docx(self, file_path: Path) -> str:
        """Extract text from Word document"""
        if DocxDocument is None:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")
        
        doc = DocxDocument(str(file_path))
        text = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(paragraph.text)
        
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text.append(row_text)
        
        return "\n\n".join(text)
    
    def _extract_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint presentation"""
        if Presentation is None:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
        
        prs = Presentation(str(file_path))
        text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"Slide {slide_num}:"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            text.append("\n".join(slide_text))
        
        return "\n\n".join(text)
    
    def _extract_text(self, file_path: Path) -> str:
        """Extract text from plain text or markdown file"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    
    def _extract_excel(self, file_path: Path) -> str:
        """Extract text from Excel spreadsheet"""
        if openpyxl is None:
            raise ImportError("openpyxl not installed. Install with: pip install openpyxl")
        
        wb = openpyxl.load_workbook(str(file_path), data_only=True)
        text = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text.append(f"Sheet: {sheet_name}")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip():
                    text.append(row_text)
        
        return "\n\n".join(text)
    
    def _compute_file_hash(self, file_path: Path) -> str:
        """Compute SHA256 hash of file for deduplication"""
        sha256_hash = hashlib.sha256()
        
        with open(file_path, "rb") as f:
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256_hash.update(byte_block)
        
        return sha256_hash.hexdigest()
    
    def get_processing_summary(self, stats: Dict[str, Any]) -> str:
        """Generate human-readable summary of processing results"""
        summary = f"""
Document Processing Summary
{'=' * 50}

Total Files Found: {stats['total_files']}
Successfully Processed: {stats['processed']}
Failed: {stats['failed']}
Skipped: {stats['skipped']}
Total Chunks Created: {stats['total_chunks']}

Files by Type:
"""
        for file_type, count in stats['files_by_type'].items():
            summary += f"  - {file_type}: {count}\n"
        
        if stats['errors']:
            summary += f"\nErrors ({len(stats['errors'])}):\n"
            for error in stats['errors'][:5]:  # Show first 5 errors
                summary += f"  - {error['file']}: {error['error']}\n"
            
            if len(stats['errors']) > 5:
                summary += f"  ... and {len(stats['errors']) - 5} more\n"
        
        return summary
